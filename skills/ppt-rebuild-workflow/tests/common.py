from __future__ import annotations

import shutil
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Inches

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def make_basic_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.name = "body-text-main"
    box.text = "测试正文"
    presentation.save(path)
    return path


def rewrite_zip(path: Path, replacements: dict[str, bytes]) -> None:
    temporary = path.with_suffix(".rewritten.pptx")
    with zipfile.ZipFile(path) as source, zipfile.ZipFile(
        temporary, "w", zipfile.ZIP_DEFLATED
    ) as target:
        for item in source.infolist():
            target.writestr(item, replacements.get(item.filename, source.read(item.filename)))
    shutil.move(temporary, path)


def make_slide_layout_relationship_absolute(path: Path) -> None:
    relationship_name = "ppt/slides/_rels/slide1.xml.rels"
    with zipfile.ZipFile(path) as archive:
        root = ET.fromstring(archive.read(relationship_name))
    for relationship in root:
        target = relationship.attrib.get("Target", "")
        if "slideLayouts/" in target:
            relationship.attrib["Target"] = "/ppt/slideLayouts/" + target.rsplit("/", 1)[-1]
    rewrite_zip(
        path,
        {relationship_name: ET.tostring(root, encoding="utf-8", xml_declaration=True)},
    )


def add_explicit_run_fonts(
    path: Path,
    latin: str,
    east_asian: str,
    complex_script: str = "ComplexFont",
    symbol: str = "SymbolFont",
) -> None:
    with zipfile.ZipFile(path) as archive:
        root = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    run_properties = root.find(f".//{{{A_NS}}}rPr")
    if run_properties is None:
        run = root.find(f".//{{{A_NS}}}r")
        if run is None:
            raise AssertionError("fixture has no text run")
        run_properties = ET.Element(f"{{{A_NS}}}rPr")
        run.insert(0, run_properties)
    for local_name, typeface in (
        ("latin", latin),
        ("ea", east_asian),
        ("cs", complex_script),
        ("sym", symbol),
    ):
        node = ET.SubElement(run_properties, f"{{{A_NS}}}{local_name}")
        node.set("typeface", typeface)
    rewrite_zip(
        path,
        {"ppt/slides/slide1.xml": ET.tostring(root, encoding="utf-8", xml_declaration=True)},
    )


def make_full_slide_picture_pptx(path: Path, image_path: Path) -> Path:
    Image.new("RGB", (1280, 720), "#284060").save(image_path)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    picture = slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    picture.name = "background-main"
    body = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.5))
    body.name = "body-text-main"
    body.text = "少量正文"
    page_number = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(0.5), Inches(0.3))
    page_number.name = "page-number-main"
    page_number.text = "01"
    presentation.save(path)
    return path


def make_connector_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(5), Inches(1))
    text.name = "body-text-main"
    text.text = "这是一段被原生连接线穿过的正文文本，用于验证碰撞审计。"
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(1),
        Inches(3),
        Inches(8),
        Inches(3),
    )
    connector.name = "decor-line-main"
    presentation.save(path)
    return path


def make_inherited_placeholder_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "继承布局标题"
    slide.placeholders[1].text = "继承布局副标题"
    presentation.save(path)
    return path


def make_rotation_and_group_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(5), Inches(1))
    text.name = "body-text-main"
    text.text = "旋转和组合对象风险测试"
    rotated = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1),
        Inches(2.45),
        Inches(7),
        Inches(0.08),
    )
    rotated.name = "decor-line-rotated"
    rotated.rotation = 5
    group = slide.shapes.add_group_shape()
    group.name = "group-main"
    grouped_line = group.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1),
        Inches(3),
        Inches(4),
        Inches(0.05),
    )
    grouped_line.name = "decor-line-grouped"
    presentation.save(path)
    return path


def make_role_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for index, name in enumerate(
        ("title-main", "body-text-main", "page-number-main", "body-panel-main")
    ):
        shape = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.5 + index),
            Inches(3),
            Inches(0.5),
        )
        shape.name = name
        shape.text = name
    image_path = path.with_suffix(".png")
    Image.new("RGB", (64, 64), "#335577").save(image_path)
    picture = slide.shapes.add_picture(str(image_path), Inches(8), Inches(1))
    picture.name = "content-image-01"
    presentation.save(path)
    return path


def make_grouped_full_slide_picture_pptx(path: Path, image_path: Path) -> Path:
    Image.new("RGB", (1280, 720), "#334455").save(image_path)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    group.name = "group-background"
    picture = group.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    picture.name = "background-main"
    presentation.save(path)
    return path
