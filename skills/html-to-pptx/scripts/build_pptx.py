#!/usr/bin/env python3
"""M3 · 构建器：layout-spec.json → .pptx（python-pptx + lxml OXML helper）。

全部绝对定位；1280×720 px 恰好满幅（px × 9525 = EMU）。baked 角色放透明占位并登记
pendingBake，交给 driver/M4 用 Playwright 补截。见 references/pipeline-contracts.md。
"""

from __future__ import annotations

import argparse
import io
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from PIL import Image

sys.path.insert(0, str(Path(__file__).resolve().parent))
import _oxml_helpers as ox  # noqa: E402

EMU_PER_PX = 9525
SLIDE_W = 12192000
SLIDE_H = 6858000
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def emu(px: float) -> int:
    return int(round(px * EMU_PER_PX))


def _rgb(hex8: str) -> RGBColor:
    h = hex8.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


_PLACEHOLDER_PNG = None


def placeholder_png() -> io.BytesIO:
    global _PLACEHOLDER_PNG
    if _PLACEHOLDER_PNG is None:
        buf = io.BytesIO()
        Image.new("RGBA", (2, 2), (0, 0, 0, 0)).save(buf, "PNG")
        _PLACEHOLDER_PNG = buf.getvalue()
    return io.BytesIO(_PLACEHOLDER_PNG)


def _has_cjk(text: str) -> bool:
    return any("一" <= c <= "鿿" or "　" <= c <= "ヿ" for c in text)


def _apply_fill(shape, fill):
    if not fill:
        shape.fill.background()
        return
    if fill["type"] == "solid":
        ox.set_solid_fill(shape, fill["color"])
    elif fill["type"] in ("linear", "radial"):
        ox.set_gradient_fill(shape, fill["stops"], fill.get("angle", 90))


def _apply_line(shape, line):
    if not line:
        shape.line.fill.background()
        return
    shape.line.color.rgb = _rgb(line["color"])
    shape.line.width = Emu(emu(line["width"]))
    dash = {"dash": "dash", "dot": "sysDot", "solid": "solid"}.get(line.get("dash"), "solid")
    if dash != "solid":
        ln = shape.line._get_or_add_ln()
        for el in ln.findall(qn("a:prstDash")):
            ln.remove(el)
        ln.append(parse_xml(f'<a:prstDash xmlns:a="{A_NS}" val="{dash}"/>'))


def build_shape(slide, sh):
    b = sh["bboxPx"]
    left, top, w, h = emu(b["x"]), emu(b["y"]), emu(b["w"]), emu(b["h"])
    is_round = bool(sh.get("radius"))
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if is_round else MSO_SHAPE.RECTANGLE, left, top, w, h
    )
    shp.shadow.inherit = False
    if is_round:
        r = sh["radius"]
        ox.set_rounded_adj(shp, min(r["tl"], r["tr"], r["br"], r["bl"]), min(b["w"], b["h"]))
    _apply_fill(shp, sh.get("fill"))
    _apply_line(shp, sh.get("line"))
    if sh.get("shadow"):
        s = sh["shadow"]
        ox.set_outer_shadow(shp, s["offsetX"], s["offsetY"], s["blur"], s["color"])
    if sh.get("rot"):
        shp.rotation = sh["rot"]
    return shp


def build_text(slide, sh):
    b = sh["bboxPx"]
    box = slide.shapes.add_textbox(emu(b["x"]), emu(b["y"]), emu(b["w"]), emu(b["h"]))
    tf = box.text_frame
    tf.word_wrap = True
    ox.disable_autofit(tf)
    t = sh["text"]
    pad = t["paddingPx"]
    tf.margin_left = Emu(emu(pad["left"]))
    tf.margin_top = Emu(emu(pad["top"]))
    tf.margin_right = Emu(emu(pad["right"]))
    tf.margin_bottom = Emu(emu(pad["bottom"]))
    tf.vertical_anchor = ANCHOR.get(t.get("valign", "top"), MSO_ANCHOR.TOP)
    para = tf.paragraphs[0]
    para.alignment = ALIGN.get(t.get("align", "left"), PP_ALIGN.LEFT)
    for run in t["runs"]:
        r = para.add_run()
        r.text = run["text"]
        r.font.bold = run["bold"]
        r.font.italic = run["italic"]
        r.font.size = Pt(run["sizePt"])
        r.font.name = run["font"]
        r.font.color.rgb = _rgb(run["color"])
        if _has_cjk(run["text"]):
            rPr = r._r.get_or_add_rPr()
            for tag in ("a:latin", "a:ea"):
                for el in rPr.findall(qn(tag)):
                    rPr.remove(el)
            rPr.append(parse_xml(f'<a:latin xmlns:a="{A_NS}" typeface="{run["font"]}"/>'))
            rPr.append(parse_xml(f'<a:ea xmlns:a="{A_NS}" typeface="{run["font"]}"/>'))
    if sh.get("fill"):
        _apply_fill(box, sh["fill"])
    if sh.get("line"):
        _apply_line(box, sh["line"])
    return box


def build_image(slide, sh, base_dir, pending):
    b = sh["bboxPx"]
    left, top, w, h = emu(b["x"]), emu(b["y"]), emu(b["w"]), emu(b["h"])
    src = (sh.get("image") or {}).get("src") or ""
    path = None
    if src.startswith("file://"):
        from urllib.parse import urlparse, unquote
        path = Path(unquote(urlparse(src).path.lstrip("/")))
    elif src and not src.startswith(("<", "data:", "http")):
        candidate = Path(src)
        if not candidate.is_absolute():
            candidate = base_dir / candidate
        path = candidate
    if path and path.exists():
        return slide.shapes.add_picture(str(path), left, top, w, h)
    pending.append({"page": sh.get("_page"), "id": sh["id"], "bboxPx": b,
                    "selectorPath": (sh.get("pendingBake") or {}).get("selectorPath", ""),
                    "targetPng": (sh.get("pendingBake") or {}).get("targetPng", "")})
    return slide.shapes.add_picture(placeholder_png(), left, top, w, h)


_GROUP_COUNTER = [1000]


def wrap_groups(slide, spec, id_to_element):
    sp_tree = slide.shapes._spTree
    for group in spec.get("groups", []):
        children = [id_to_element[cid] for cid in group["children"] if cid in id_to_element]
        if len(children) < 2:
            continue
        xs, ys, x2s, y2s = [], [], [], []
        for cid in group["children"]:
            sh = spec["_shape_by_id"].get(cid)
            if not sh:
                continue
            b = sh["bboxPx"]
            xs.append(b["x"]); ys.append(b["y"]); x2s.append(b["x"] + b["w"]); y2s.append(b["y"] + b["h"])
        if not xs:
            continue
        ox_, oy = emu(min(xs)), emu(min(ys))
        ex, ey = emu(max(x2s)) - ox_, emu(max(y2s)) - oy
        _GROUP_COUNTER[0] += 1
        grp = parse_xml(
            f'<p:grpSp xmlns:p="{P_NS}" xmlns:a="{A_NS}">'
            f'<p:nvGrpSpPr><p:cNvPr id="{_GROUP_COUNTER[0]}" name="group-{_GROUP_COUNTER[0]}"/>'
            f'<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            f'<p:grpSpPr><a:xfrm><a:off x="{ox_}" y="{oy}"/><a:ext cx="{ex}" cy="{ey}"/>'
            f'<a:chOff x="{ox_}" y="{oy}"/><a:chExt cx="{ex}" cy="{ey}"/></a:xfrm></p:grpSpPr>'
            f'</p:grpSp>'
        )
        for el in children:
            sp_tree.remove(el)
            grp.append(el)
        sp_tree.append(grp)


def build_slide(prs, spec, base_dir, pending):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec["_shape_by_id"] = {sh["id"]: sh for sh in spec["shapes"]}
    id_to_element = {}
    for sh in sorted(spec["shapes"], key=lambda s: s["zOrder"]):
        sh["_page"] = spec.get("page")
        role = sh["role"]
        baked = sh.get("expressibility", {}).get("verdict") == "baked"
        if baked:
            shp = build_image(slide, sh, base_dir, pending)  # 烘焙元素统一放 picture
        elif role == "text":
            shp = build_text(slide, sh)
        elif role in ("image", "svg"):
            shp = build_image(slide, sh, base_dir, pending)
        elif role == "table":
            continue  # v1：table 由其单元格 text 形状呈现，原生表格留 v2
        else:
            shp = build_shape(slide, sh)
        id_to_element[sh["id"]] = shp._element
    wrap_groups(slide, spec, id_to_element)
    return slide


def build(spec_paths: list[Path], output: Path) -> dict:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    specs = []
    for p in spec_paths:
        specs.append((p, json.loads(p.read_text(encoding="utf-8"))))
    specs.sort(key=lambda item: item[1].get("page", 0))
    pending: list[dict] = []
    for path, spec in specs:
        build_slide(prs, spec, path.parent, pending)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    bake_path = output.with_suffix(output.suffix + ".pending-bake.json")
    bake_path.write_text(json.dumps({"pendingBake": pending}, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"pptx": str(output), "pendingBake": len(pending), "bakeManifest": str(bake_path)}


def make_stdout_robust() -> None:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(errors="backslashreplace")


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("specs", nargs="+", help="layout-spec.json path(s)")
    parser.add_argument("--output", required=True)
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args()
    try:
        result = build([Path(s) for s in args.specs], Path(args.output))
    except (OSError, json.JSONDecodeError) as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 2
    make_stdout_robust()
    print(json.dumps(result, ensure_ascii=False) if args.json else result["pptx"])
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
