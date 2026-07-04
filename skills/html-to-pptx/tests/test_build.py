from __future__ import annotations

import json
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation

SKILL_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(SKILL_ROOT / "scripts"))

import build_pptx  # noqa: E402

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS = {"a": A, "p": P}


def shape(sid, role, x, y, w, h, z=0, **kw):
    s = {
        "id": sid, "role": role, "bboxPx": {"x": x, "y": y, "w": w, "h": h},
        "rot": 0.0, "zOrder": z, "groupId": None,
        "fill": None, "line": None, "shadow": None, "radius": None,
        "text": None, "image": None, "table": None,
        "expressibility": {"verdict": "native", "bakedReason": None}, "pendingBake": None,
    }
    s.update(kw)
    return s


def text_shape(sid, x, y, w, h, content="Hi", font="Arial", size=15.0, z=0, **kw):
    s = shape(sid, "text", x, y, w, h, z, **kw)
    s["text"] = {
        "runs": [{"text": content, "bold": False, "italic": False, "color": "#112233FF", "font": font, "sizePt": size}],
        "align": "left", "valign": "top", "lineBreaks": [],
        "paddingPx": {"left": 4, "top": 4, "right": 4, "bottom": 4},
        "textLayoutBudget": {"lines": 1, "lineWidths": [w - 8]},
    }
    return s


def build(tmp_path, shapes, groups=None, page=1):
    spec = {
        "schemaVersion": "1.0", "page": page, "mode": "html-to-pptx",
        "coordinateSystem": {"unit": "csspx", "width": 1280, "height": 720, "emuPerPx": 9525},
        "reference": None, "cleanliness": {"emittedShapes": len(shapes), "visibleElements": len(shapes), "ratio": 1.0},
        "fontMap": [], "warnings": [], "shapes": shapes, "groups": groups or [],
    }
    spec_path = tmp_path / f"spec-{page}.json"
    spec_path.write_text(json.dumps(spec), encoding="utf-8")
    out = tmp_path / "out.pptx"
    result = build_pptx.build([spec_path], out)
    return out, result


def slide_xml(pptx_path) -> ET.Element:
    with zipfile.ZipFile(pptx_path) as zf:
        return ET.fromstring(zf.read("ppt/slides/slide1.xml"))


def test_solid_round_shadow_rotation(tmp_path):
    s = shape("s1", "shape", 100, 100, 200, 120, radius={"tl": 16, "tr": 16, "br": 16, "bl": 16},
              fill={"type": "solid", "color": "#ff8800ff"},
              shadow={"offsetX": 0, "offsetY": 8, "blur": 24, "color": "#00000040"}, rot=30.0)
    out, _ = build(tmp_path, [s])
    prs = Presentation(str(out))
    shp = list(prs.slides[0].shapes)[0]
    assert shp.rotation == 30.0
    root = slide_xml(out)
    sp = root.find(".//p:sp", NS)
    assert sp.find(".//a:solidFill/a:srgbClr", NS).get("val") == "FF8800"
    assert sp.find(".//a:prstGeom", NS).get("prst") == "roundRect"
    assert sp.find(".//a:avLst/a:gd", NS) is not None
    assert sp.find(".//a:effectLst/a:outerShdw", NS) is not None


def test_gradient_stops_written(tmp_path):
    s = shape("g1", "shape", 0, 0, 400, 200,
              fill={"type": "linear", "stops": [{"color": "#ff0000ff", "offset": 0.0}, {"color": "#0000ffff", "offset": 1.0}], "angle": 90})
    out, _ = build(tmp_path, [s])
    root = slide_xml(out)
    gs = root.findall(".//a:gradFill/a:gsLst/a:gs", NS)
    assert len(gs) == 2
    assert gs[0].find("a:srgbClr", NS).get("val") == "FF0000"
    assert gs[0].get("pos") == "0"
    assert gs[1].get("pos") == "100000"
    assert root.find(".//a:gradFill/a:lin", NS).get("ang") == "0"  # css 90 -> ooxml 0


def test_textbox_no_autofit_and_run_props(tmp_path):
    out, _ = build(tmp_path, [text_shape("t1", 50, 50, 300, 60, content="Hello", font="Segoe UI", size=18.0)])
    root = slide_xml(out)
    tx = root.find(".//p:sp/p:txBody", NS)
    assert tx.find("a:bodyPr/a:noAutofit", NS) is not None
    run = tx.find(".//a:r", NS)
    assert run.find("a:t", NS).text == "Hello"
    rpr = run.find("a:rPr", NS)
    assert rpr.get("sz") == "1800"  # 18pt * 100
    assert rpr.find("a:latin", NS).get("typeface") == "Segoe UI"
    assert rpr.find("a:solidFill/a:srgbClr", NS).get("val") == "112233"


def test_cjk_sets_eastasian_typeface(tmp_path):
    out, _ = build(tmp_path, [text_shape("t1", 0, 0, 300, 60, content="标题文字", font="Microsoft YaHei", size=20.0)])
    root = slide_xml(out)
    rpr = root.find(".//a:r/a:rPr", NS)
    assert rpr.find("a:latin", NS).get("typeface") == "Microsoft YaHei"
    assert rpr.find("a:ea", NS).get("typeface") == "Microsoft YaHei"


def test_group_wraps_children_with_identity_transform(tmp_path):
    a = text_shape("a", 10, 10, 100, 30, content="A", z=0, groupId="g")
    b = text_shape("b", 10, 50, 100, 30, content="B", z=1, groupId="g")
    groups = [{"id": "g", "children": ["a", "b"], "label": "card"}]
    out, _ = build(tmp_path, [a, b], groups)
    root = slide_xml(out)
    grp = root.find(".//p:grpSp", NS)
    assert grp is not None
    assert len(grp.findall("p:sp", NS)) == 2
    xfrm = grp.find("p:grpSpPr/a:xfrm", NS)
    off, ext = xfrm.find("a:off", NS), xfrm.find("a:ext", NS)
    choff, chext = xfrm.find("a:chOff", NS), xfrm.find("a:chExt", NS)
    assert (off.get("x"), off.get("y")) == (choff.get("x"), choff.get("y"))
    assert (ext.get("cx"), ext.get("cy")) == (chext.get("cx"), chext.get("cy"))


def test_baked_image_uses_placeholder_and_registers_pending(tmp_path):
    s = shape("img1", "image", 0, 0, 200, 200,
              image={"src": "<svg></svg>", "objectFit": "cover", "srcRect": None},
              expressibility={"verdict": "baked", "bakedReason": "svg"},
              pendingBake={"selectorPath": "html > body > svg", "targetPng": "bake/page-1-svg.png"})
    out, result = build(tmp_path, [s])
    assert result["pendingBake"] == 1
    manifest = json.loads(Path(result["bakeManifest"]).read_text(encoding="utf-8"))
    assert manifest["pendingBake"][0]["id"] == "img1"
    root = slide_xml(out)
    assert root.find(".//p:pic", NS) is not None


def test_shape_and_text_counts_match_spec(tmp_path):
    shapes = [shape("bg", "shape", 0, 0, 1280, 720, z=0, fill={"type": "solid", "color": "#ffffffff"}),
              text_shape("t1", 100, 100, 300, 40, z=1),
              text_shape("t2", 100, 200, 300, 40, z=2)]
    out, _ = build(tmp_path, shapes)
    prs = Presentation(str(out))
    kinds = [s.shape_type for s in prs.slides[0].shapes]
    assert len(kinds) == 3
